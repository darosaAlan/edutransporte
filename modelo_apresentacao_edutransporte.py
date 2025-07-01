from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Cria a apresentação
prs = Presentation()

# Slide de capa
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
if slide.shapes.title is not None:
    slide.shapes.title.text_frame.text = "EduTransporte"
slide.placeholders[1].text = "Apresentação das Telas do App\nAluno e Motorista"

# Slide do fluxograma
diagram_slide = prs.slides.add_slide(prs.slide_layouts[5])
if diagram_slide.shapes.title is not None:
    diagram_slide.shapes.title.text_frame.text = "Fluxograma de Navegação"
# Adiciona caixa de texto manualmente
left = Inches(0.5)
top = Inches(1.5)
width = Inches(8)
height = Inches(2)
txBox = diagram_slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
tf.text = "(Insira aqui o fluxograma do app)"
tf.paragraphs[0].font.size = Pt(20)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Dados das telas
telas = [
    ("Tela de Login", "Permite login de aluno ou motorista, com seleção de perfil."),
    ("Tela de Cadastro", "Cadastro de novo usuário, validação de dados."),
    ("Dashboard do Aluno", "Acesso rápido à localização da van, mensagens, agendamento, configurações e logout."),
    ("Dashboard do Motorista", "Gerenciamento de viagem, localização, mensagens, configurações e logout."),
    ("Tela de Localização", "Mapa em tempo real com posição da van."),
    ("Tela de Mensagens", "Chat entre aluno e motorista."),
    ("Tela de Agendamento", "Agendamento de viagens e confirmação de presença."),
    ("Tela de Configurações", "Alteração de dados do perfil.")
]

for nome, desc in telas:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    if slide.shapes.title is not None:
        slide.shapes.title.text_frame.text = nome
    # Espaço reservado para imagem
    left = Inches(0.5)
    top = Inches(1.2)
    width = Inches(2.5)
    height = Inches(4.5)
    img_placeholder = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    if hasattr(img_placeholder, 'text_frame'):
        img_placeholder.text_frame.text = "(Insira aqui o print da tela)"
        img_placeholder.text_frame.paragraphs[0].font.size = Pt(14)
        img_placeholder.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    # Descrição
    txBox = slide.shapes.add_textbox(Inches(3.2), Inches(1.2), Inches(5), Inches(2))
    tf = txBox.text_frame
    tf.text = desc
    tf.paragraphs[0].font.size = Pt(18)

# Salva o arquivo
prs.save("Apresentacao_EduTransporte.pptx")
print("Apresentação gerada com sucesso: Apresentacao_EduTransporte.pptx") 